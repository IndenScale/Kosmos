import os
import tempfile
import uuid
from typing import Tuple, List
from pathlib import Path
from .base_processor import BaseProcessor
import pypandoc
import zipfile
from PIL import Image
import shutil
from pptx import Presentation
import xml.etree.ElementTree as ET

class PptxProcessor(BaseProcessor):
    """PPTX文档处理器"""
    
    def __init__(self):
        super().__init__()
        self.supported_extensions = ['.pptx']
        self.temp_dir = Path("temp/pptx_processing")
        self.temp_dir.mkdir(parents=True, exist_ok=True)
    
    def can_process(self, file_path: str) -> bool:
        """判断是否可以处理该文件类型"""
        file_ext = Path(file_path).suffix.lower()
        return file_ext in self.supported_extensions
    
    def _extract_content_impl(self, file_path: str) -> Tuple[str, List[str]]:
        """提取PPTX文档内容和图片，按页面分割"""
        try:
            # 创建临时目录
            temp_id = str(uuid.uuid4())
            temp_extract_dir = self.temp_dir / temp_id
            temp_extract_dir.mkdir(exist_ok=True)
            
            try:
                # 1. 提取图片
                image_paths = self._extract_images(file_path, temp_extract_dir)
                
                # 2. 按页面提取内容
                slides_content = self._extract_slides_content(file_path)
                
                # 3. 将每个页面转换为markdown格式
                markdown_text = self._convert_slides_to_markdown(slides_content)
                
                # 4. 处理图片引用
                markdown_text = self._process_image_references(markdown_text, image_paths)
                
                return markdown_text, image_paths
                
            finally:
                # 清理临时目录
                if temp_extract_dir.exists():
                    shutil.rmtree(temp_extract_dir)
                    
        except Exception as e:
            raise Exception(f"PPTX处理器提取内容失败: {str(e)}")
    
    def _extract_slides_content(self, pptx_path: str) -> List[dict]:
        """提取每个幻灯片的内容"""
        slides_content = []
        
        try:
            prs = Presentation(pptx_path)
            
            for slide_num, slide in enumerate(prs.slides, 1):
                slide_data = {
                    'slide_number': slide_num,
                    'title': '',
                    'content': [],
                    'notes': ''
                }
                
                # 提取幻灯片内容
                for shape in slide.shapes:
                    if hasattr(shape, "text") and shape.text.strip():
                        # 判断是否为标题
                        if shape.is_placeholder and shape.placeholder_format.type == 1:  # 标题占位符
                            slide_data['title'] = shape.text.strip()
                        else:
                            slide_data['content'].append(shape.text.strip())
                
                # 提取备注
                if slide.notes_slide and slide.notes_slide.notes_text_frame:
                    slide_data['notes'] = slide.notes_slide.notes_text_frame.text.strip()
                
                slides_content.append(slide_data)
                
        except Exception as e:
            print(f"提取幻灯片内容失败: {str(e)}")
            # 如果python-pptx失败，尝试使用pypandoc作为后备
            try:
                markdown_text = pypandoc.convert_file(pptx_path, 'markdown')
                # 简单按分隔符分割页面
                slides = markdown_text.split('\n---\n')
                for i, slide_content in enumerate(slides, 1):
                    slides_content.append({
                        'slide_number': i,
                        'title': f'幻灯片 {i}',
                        'content': [slide_content.strip()],
                        'notes': ''
                    })
            except Exception as fallback_e:
                print(f"后备方案也失败: {str(fallback_e)}")
        
        return slides_content
    
    def _convert_slides_to_markdown(self, slides_content: List[dict]) -> str:
        """将幻灯片内容转换为markdown格式，每个页面作为一个chunk"""
        markdown_parts = []
        
        for slide_data in slides_content:
            slide_md = []
            
            # 添加页面分隔符和页码
            slide_md.append(f"\n\n--- 幻灯片 {slide_data['slide_number']} ---\n")
            
            # 添加标题
            if slide_data['title']:
                slide_md.append(f"# {slide_data['title']}\n")
            
            # 添加内容
            for content in slide_data['content']:
                if content:
                    slide_md.append(f"{content}\n")
            
            # 添加备注
            if slide_data['notes']:
                slide_md.append(f"\n**备注:** {slide_data['notes']}\n")
            
            # 添加页面结束标记
            slide_md.append("\n--- 页面结束 ---\n")
            
            markdown_parts.append(''.join(slide_md))
        
        return '\n'.join(markdown_parts)
    
    def _extract_images(self, pptx_path: str, extract_dir: Path) -> List[str]:
        """从PPTX文件中提取图片"""
        image_paths = []
        
        try:
            with zipfile.ZipFile(pptx_path, 'r') as pptx_zip:
                # 查找图片文件
                for file_info in pptx_zip.filelist:
                    if file_info.filename.startswith('ppt/media/'):
                        # 提取图片
                        image_data = pptx_zip.read(file_info.filename)
                        
                        # 生成图片文件名
                        image_name = Path(file_info.filename).name
                        image_path = extract_dir / image_name
                        
                        # 保存图片
                        with open(image_path, 'wb') as img_file:
                            img_file.write(image_data)
                        
                        image_paths.append(str(image_path))
                        
        except Exception as e:
            print(f"提取图片失败: {str(e)}")
        
        return image_paths
    
    def _process_image_references(self, markdown_text: str, image_paths: List[str]) -> str:
        """处理markdown中的图片引用"""
        for image_path in image_paths:
            image_name = Path(image_path).name
            # 在markdown中查找并替换图片引用
            markdown_text = markdown_text.replace(
                f"![](media/{image_name})",
                f"![{image_name}](media/{image_name})\n\n[图片描述占位符: {image_path}]\n"
            )
        
        return markdown_text
    
    def get_supported_extensions(self) -> List[str]:
        return self.supported_extensions